from pptx import Presentation
from pptx.util import Inches
import os

def create_ppt(images_folder, output_pptx):
    # Create a presentation object
    presentation = Presentation()

    # Get a list of image files in the specified folder
    image_files = [f for f in os.listdir(images_folder) if f.endswith(('.png', '.jpg', '.jpeg', '.gif'))]

    # Sort the image files based on their names
    # image_files.sort()
    image_files = sorted(image_files, key=lambda x: int(x.split('_')[1].split('.')[0]))
    # Iterate through the image files and add each image to a new slide
    for image_file in image_files:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout

        image_path = os.path.join(images_folder, image_file)
        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(9))

    # Save the presentation to a file
    presentation.save(output_pptx)

if __name__ == "__main__":

    # Specify the folder containing images and the output PowerPoint file
    images_folder = 'images'
    output_pptx = 'presentation1.pptx'

    # Create the PowerPoint presentation
    create_ppt(images_folder, output_pptx)

    print(f"PowerPoint presentation created at: {output_pptx}")
